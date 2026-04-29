"""Build ActivitiesApp final presentation PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

W = Inches(13.333)
H = Inches(7.5)

BG      = RGBColor(0x0D, 0x0D, 0x0D)   # near black
ACCENT  = RGBColor(0x4A, 0x90, 0xE2)   # blue
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xAA, 0xAA, 0xAA)
GREEN   = RGBColor(0x4C, 0xAF, 0x50)
YELLOW  = RGBColor(0xFF, 0xC1, 0x07)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank


def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def box(slide, left, top, width, height, text="", size=18, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True, fill=None, fill_color=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if fill and fill_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fill_color
    return txBox


def heading(slide, text, sub=None):
    box(slide, 0.5, 0.25, 12.3, 1.0, text, size=36, bold=True,
        color=ACCENT, align=PP_ALIGN.LEFT)
    if sub:
        box(slide, 0.5, 1.1, 12.3, 0.5, sub, size=18, color=GRAY, align=PP_ALIGN.LEFT)


def hline(slide, top):
    line = slide.shapes.add_connector(
        1,  # straight
        Inches(0.5), Inches(top), Inches(12.8), Inches(top))
    line.line.color.rgb = ACCENT
    line.line.width = Pt(1)


def bullet_block(slide, left, top, width, height, items, size=16, color=WHITE, title=None, title_color=None):
    if title:
        box(slide, left, top, width, 0.4, title, size=18, bold=True,
            color=title_color or ACCENT)
        top += 0.42
        height -= 0.42
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = ("• " if not item.startswith("  ") else "    ◦ ") + item.lstrip()
        run.font.size = Pt(size)
        run.font.color.rgb = color


def col_card(slide, left, top, width, height, title, items, t_color=None):
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)
    box(slide, left + 0.12, top + 0.1, width - 0.24, 0.38,
        title, size=16, bold=True, color=t_color or ACCENT)
    bullet_block(slide, left + 0.12, top + 0.5, width - 0.24,
                 height - 0.6, items, size=13)


# ─── SLIDE 1: Title ───────────────────────────────────────────────────────────
s = add_slide()
box(s, 1, 1.5, 11.3, 1.5, "ActivitiesApp", size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 1, 3.1, 11.3, 0.7,
    "SE 3630 / SE 3830 — Final Presentation",
    size=24, color=ACCENT, align=PP_ALIGN.CENTER)
box(s, 1, 3.9, 11.3, 0.5,
    "Samuel Schuetz  ·  Cael Church",
    size=20, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 1, 5.0, 11.3, 0.5,
    "Discover local activities — nearby search, user events, maps, web + mobile",
    size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ─── SLIDE 2: SE 3630 Learning Outcomes ───────────────────────────────────────
s = add_slide()
heading(s, "SE 3630 — Mobile App Learning Outcomes")
hline(s, 1.5)

outcomes = [
    ("Mobile Dev Environment",
     ["Built with .NET MAUI — one codebase, Android + iOS + Windows",
      "Shared Razor pages via ActivitiesApp.Shared (Razor class library)",
      "[screenshot: src/ActivitiesApp.Maui/MauiProgram.cs]"],
     GREEN),
    ("App Lifecycle",
     ["MauiNetworkStatus hooks into IConnectivity.ConnectivityChanged",
      "Offline cache syncs on reconnect via delta sync (PullChanges / PushChanges gRPC)",
      "[screenshot: src/ActivitiesApp.Maui/Services/MauiNetworkStatus.cs]"],
     ACCENT),
    ("UI Layout & Design",
     ["Shared Razor pages: Home, Activities, ActivityDetail, CreateActivity, Profile",
      "Location popup, category pills, fuzzy search bar — all in Blazor Hybrid",
      "[screenshot: src/ActivitiesApp.Shared/Pages/Home.razor lines 16–60]"],
     YELLOW),
    ("Data Binding",
     ["Blazor two-way binding (@bind) for zip input, search, filters",
      "ActivitiesViewModel inherits BaseViewModel; state drives UI automatically",
      "[screenshot: src/ActivitiesApp.Maui/ViewModels/ActivitiesViewModel.cs]"],
     RGBColor(0xFF, 0x70, 0x50)),
]

col_w = 3.0
gap = 0.1
start_l = 0.5
for i, (title, items, color) in enumerate(outcomes):
    col_card(s, start_l + i * (col_w + gap), 1.65, col_w, 5.5, title, items, color)

# ─── SLIDE 3: SE 3630 Outcomes cont. ─────────────────────────────────────────
s = add_slide()
heading(s, "SE 3630 — GUI Patterns & Testing")
hline(s, 1.5)

col_card(s, 0.5, 1.65, 5.9, 5.5,
    "GUI Framework Patterns",
    ["MVVM: ViewModels own state, Razor pages observe",
     "Dependency injection — IActivityService, INetworkStatus, ILocationProvider",
     "NavigationManager for push / pop between pages",
     "Blazor Hybrid = WebView wrapping native Razor rendering",
     "[screenshot: src/ActivitiesApp.Shared/Pages/Activities.razor]"],
    ACCENT)

col_card(s, 6.6, 1.65, 6.2, 5.5,
    "Testing Mobile Code",
    ["Unit tests: HelpersTests, FuzzySearchServiceTests, ActivityFilterServiceTests",
     "ViewModels tested in ActivitiesApp.ViewModels.Tests",
     "Integration tests hit real Postgres (no mocks) — catches boundary bugs",
     "Playwright E2E tests in playwright-tests/",
     "CI runs lint → unit → integration → build on every PR",
     "[screenshot: tests/ActivitiesApp.Core.Tests/HelpersTests.cs]"],
    GREEN)

# ─── SLIDE 4: SE 3830 Learning Outcomes ───────────────────────────────────────
s = add_slide()
heading(s, "SE 3830 — Cloud Learning Outcomes")
hline(s, 1.5)

cloud = [
    ("Deploy to Public Cloud",
     ["App runs on Kubernetes in the cloud at activor.duckdns.org",
      "Self-hosted GitHub Actions runner deploys on every master push",
      "[screenshot: .github/workflows/deploy.yml]"],
     GREEN),
    ("Virt vs Containerization",
     ["Docker containers package API + Web — image tagged with Git SHA",
      "K8s pods schedule containers; each pod is isolated, scalable",
      "No VMs — pure container workloads on shared cluster"],
     ACCENT),
    ("Security Controls",
     ["Secrets in K8s Secrets, not env files — injected at runtime",
      "HTTPS via TLS certs auto-renewed by weekly CronJob (certbot + DuckDNS)",
      "Microsoft Entra ID (OIDC) — no passwords stored",
      "RBAC: cert-renewer ServiceAccount scoped to only its secrets"],
     RGBColor(0xFF, 0x70, 0x50)),
    ("Cloud Models (IaaS / PaaS / SaaS)",
     ["IaaS: bare K8s cluster (we manage everything)",
      "PaaS: Azure Blob Storage (backup), GCP Monitoring (metrics)",
      "SaaS: Google Places API, Microsoft Entra ID, DuckDNS",
      "Chose IaaS K8s for control; SaaS for third-party integrations"],
     YELLOW),
]

col_w = 3.0
gap = 0.1
start_l = 0.5
for i, (title, items, color) in enumerate(cloud):
    col_card(s, start_l + i * (col_w + gap), 1.65, col_w, 5.5, title, items, color)

# ─── SLIDE 5: SE 3830 Outcomes cont. ─────────────────────────────────────────
s = add_slide()
heading(s, "SE 3830 — CI/CD & Automation")
hline(s, 1.5)

col_card(s, 0.5, 1.65, 5.9, 5.5,
    "Deploy & Manage Code in Cloud",
    ["Pipeline: lint → unit tests → integration tests → build images → deploy",
     "Docker images tagged with short Git SHA — every deploy traceable",
     "K8s rolling update with zero-downtime verified by k6 smoke test",
     "Database migrations run as K8s Job before pods swap to new image",
     "Postgres daily backups to Azure Blob Storage via CronJob",
     "[screenshot: .github/workflows/deploy.yml (deploy job)]"],
    ACCENT)

col_card(s, 6.6, 1.65, 6.2, 5.5,
    "Automate Compile / Test / Deploy",
    ["pr-environment.yml: lint + unit + integration + image build on every PR",
     "deploy.yml: full pipeline on master push — no manual steps",
     "TLS cert renewal CronJob fires weekly (Wed 14:18 MT, certbot + DuckDNS)",
     "Uptime Kuma monitors activor.duckdns.org continuously",
     "Push notification sent via ntfy if pipeline fails",
     "[screenshot: deploy/k8s/certs/cert-renew-cronjob.yaml]"],
    GREEN)

# ─── SLIDE 6: App Demo placeholder ────────────────────────────────────────────
s = add_slide()
heading(s, "App Demo")
hline(s, 1.5)

box(s, 0.5, 1.65, 12.3, 0.5,
    "Live demo — screen-share phone running ActivitiesApp on Android",
    size=18, color=GRAY)

demo_items = [
    ("Home — Discover nearby", ["GPS location detected, Google Places results + user events merged",
     "Category pills filter (Sports, Food, Music, Outdoors, …)",
     "[screenshot of web app home page: activor.duckdns.org]"]),
    ("Activities list", ["Distance-sorted cards with photos from Google Places",
     "Fuzzy search — typo tolerant (e.g. 'resturant' finds Restaurant)",
     "[screenshot of Activities page]"]),
    ("Activity Detail", ["Name, description, cost, age range, time, map pin",
     "Edit/Delete visible only to owner",
     "[screenshot of ActivityDetail page]"]),
    ("Create Activity", ["Sign in with Microsoft → form unlocks",
     "Fields: name, city, time, cost, age range, category, description",
     "[screenshot of CreateActivity page]"]),
]

for i, (title, items) in enumerate(demo_items):
    col = i % 2
    row = i // 2
    col_card(s,
             0.5 + col * 6.55, 2.3 + row * 2.5,
             6.3, 2.35,
             title, items, ACCENT)

# ─── SLIDE 7: Backing Services ─────────────────────────────────────────────────
s = add_slide()
heading(s, "Backing Services Tour")
hline(s, 1.5)

services = [
    ("gRPC API", ["ASP.NET Core — owns all data, auth, Google Places calls",
      "Proto: ListActivities (server streaming), DiscoverActivities, delta sync",
      "[screenshot: src/ActivitiesApp.Api/Protos/activities.proto lines 9–38]"],
     ACCENT),
    ("PostgreSQL 16", ["Primary data store — EF Core migrations",
      "Daily backup CronJob → Azure Blob Storage",
      "Seeded from Cosmos DB on first boot"],
     GREEN),
    ("Google Places API", ["Nearby search: GooglePlacesService queries on DiscoverActivities",
      "Photos proxied through API endpoint with 24h in-memory cache",
      "[screenshot: src/ActivitiesApp.Api/Services/GooglePlacesService.cs]"],
     YELLOW),
    ("Kubernetes Cluster", ["Manifests in deploy/k8s/ — namespace, app, data, certs, observability",
      "Ingress: activor.duckdns.org (web), grafana.activor.duckdns.org",
      "[screenshot: deploy/k8s/app/]"],
     RGBColor(0xFF, 0x70, 0x50)),
    ("Observability Stack", ["OTel Collector → Prometheus (metrics) + Loki (logs) → Grafana",
      "Panels: p95 latency, success rate, 5xx count, concurrent users, GCP API usage",
      "[screenshot of Grafana dashboard at grafana.activor.duckdns.org]"],
     RGBColor(0x9C, 0x27, 0xB0)),
    ("Microsoft Entra ID", ["OIDC — Sign in with Microsoft, no passwords stored",
      "Redirect URI: https://activor.duckdns.org/signin-oidc",
      "MauiAuthenticationStateProvider on mobile"],
     GRAY),
]

col_w = 4.05
gap = 0.12
start_l = 0.5
for i, (title, items, color) in enumerate(services):
    col = i % 3
    row = i // 3
    col_card(s,
             start_l + col * (col_w + gap),
             1.65 + row * 2.8,
             col_w, 2.65,
             title, items, color)

# ─── SLIDE 8: Code Review Improvements ────────────────────────────────────────
s = add_slide()
heading(s, "Improvements from Code Reviews")
hline(s, 1.5)

col_card(s, 0.5, 1.65, 5.9, 5.5,
    "Security & Auth Fixes",
    ["Hid Create Activity form for signed-out users — form was accessible without auth",
     "Redirect cancelled sign-ins to home instead of error page",
     "Root-relative sign-in link (was breaking on deep routes)",
     "Fixed security vulnerabilities flagged in review (commits 5b2437f, de765b7)",
     "[screenshot: src/ActivitiesApp.Shared/Pages/CreateActivity.razor — auth guard]"],
    RGBColor(0xFF, 0x70, 0x50))

col_card(s, 6.6, 1.65, 6.2, 5.5,
    "Error Handling & UX",
    ["Single error route handler — lowercase + uppercase routes both caught",
     "Errors now redirect home instead of dead-end page",
     "Zero-downtime deployment verified by improved k6 smoke test",
     "Treats compiler warnings as errors (ae9388a) — enforces code quality",
     "Pipeline now posts PR comment with environment status",
     "[screenshot: src/ActivitiesApp.Shared/Pages/NotFound.razor]"],
    ACCENT)

# ─── SLIDE 9: Individual Takeaways ────────────────────────────────────────────
s = add_slide()
heading(s, "Biggest Takeaways — Samuel & Cael")
hline(s, 1.5)

col_card(s, 0.5, 1.65, 5.9, 5.5,
    "Samuel — Top 2 Lessons",
    ["1. Kubernetes is hard until you understand the layering. \
Getting TLS, ingress, gRPC, and rolling deploys to work together \
required understanding each layer independently — no shortcut. \
The cert-renew CronJob alone took multiple debugging iterations.",

     "2. CI/CD is a forcing function for code quality. \
Having lint + unit + integration + smoke tests run on every push \
meant broken code never sat in master — but it also meant every \
shortcut eventually caught up with us."],
    WHITE)

col_card(s, 6.6, 1.65, 6.2, 5.5,
    "Cael — Top 2 Lessons",
    ["1. Observability is not optional. \
Adding Grafana + Prometheus + Loki made it obvious when something \
was wrong in production — concurrent users, 5xx rate, and log \
volume gave real signals, not guesses.",

     "2. Mobile lifecycle is different from web. \
The MAUI app needed explicit handling for offline states, network \
reconnection, and data sync — things a web app never has to think \
about. IConnectivity and the delta sync protocol were the answer."],
    WHITE)

# ─── SLIDE 10: Architecture Summary ───────────────────────────────────────────
s = add_slide()
heading(s, "Architecture at a Glance")
hline(s, 1.5)

box(s, 0.5, 1.65, 12.3, 0.4,
    "MAUI (Android/iOS/Win)   ──gRPC──▶   ASP.NET Core API   ◀──REST──   Blazor Web",
    size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
box(s, 0.5, 2.1, 12.3, 0.35,
    "│                                       │                                │",
    size=16, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 0.5, 2.45, 12.3, 0.35,
    "Offline Cache                    PostgreSQL 16                     Grafana / Prometheus / Loki",
    size=16, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 0.5, 2.8, 12.3, 0.35,
    "(delta sync)                     Cosmos DB (seed)                  OTel Collector",
    size=14, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 0.5, 3.15, 12.3, 0.35,
    "                                 Google Places API                 GCP Cloud Monitoring",
    size=14, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 0.5, 3.5, 12.3, 0.35,
    "                                 Microsoft Entra ID (OIDC)         Uptime Kuma",
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

hline(s, 4.2)

box(s, 0.5, 4.3, 12.3, 0.4, "CI/CD Pipeline (GitHub Actions)", size=18, bold=True, color=ACCENT)
box(s, 0.5, 4.75, 12.3, 0.35,
    "lint  →  unit tests  →  integration tests (real Postgres)  →  build Docker images (SHA tag)  →  deploy to K8s",
    size=15, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 0.5, 5.15, 12.3, 0.35,
    "Self-hosted runner applies K8s manifests  ·  EF migrations Job  ·  k6 smoke test  ·  ntfy alert on failure",
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

hline(s, 5.6)
box(s, 0.5, 5.7, 12.3, 0.35,
    "Live at:  activor.duckdns.org   ·   grafana.activor.duckdns.org   ·   prometheus.activor.duckdns.org",
    size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

prs.save("presentation.pptx")
print("Done — presentation.pptx written.")
